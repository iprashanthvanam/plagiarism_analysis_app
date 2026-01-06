









import os
import mimetypes
import logging
from typing import Optional

import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation

from app.libs.gemini_service import extract_text_with_gemini


logger = logging.getLogger(__name__)



async def extract_text(file_path: str, content_type: Optional[str] = None) -> str:
    """
    Unified extraction entrypoint.
    MUST return str (never None).
    """
    if not os.path.exists(file_path):
        logger.error("File does not exist: %s", file_path)
        return ""

    ext = os.path.splitext(file_path)[1].lower()
    mime, _ = mimetypes.guess_type(file_path)

    try:
        # ---------- TXT ----------
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

        # ---------- DOCX ----------
        if ext == ".docx":
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        # ---------- DOC (fallback to Gemini OCR) ----------
        if ext == ".doc":
            logger.info("Routing .doc to Gemini OCR")
            if extract_text_with_gemini:
                return await extract_text_with_gemini(file_path, is_pdf=False)
            else:
                logger.warning("Gemini disabled due to dependency conflict")
            return ""


        # ---------- XLSX ----------
        if ext == ".xlsx":
            dfs = pd.read_excel(file_path, sheet_name=None, engine="openpyxl")
            return "\n".join(
                df.astype(str).to_string(index=False) for df in dfs.values()
            )

        # ---------- XLS ----------
        if ext == ".xls":
            dfs = pd.read_excel(file_path, sheet_name=None, engine="xlrd")
            return "\n".join(
                df.astype(str).to_string(index=False) for df in dfs.values()
            )

        # ---------- PPTX ----------
        if ext == ".pptx":
            prs = Presentation(file_path)
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            return "\n".join(slides_text)

        # ---------- PPT (unsupported but safe) ----------
        if ext == ".ppt":
            logger.warning(".ppt format not supported; skipping")
            return ""

        # ---------- PDF ----------
        if ext == ".pdf":
            text = ""
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text() or ""
                    text += page_text + "\n"

            if text.strip():
                return text

            logger.info("PDF appears scanned; routing to Gemini OCR")
            if extract_text_with_gemini:
                return await extract_text_with_gemini(file_path, is_pdf=False)
            else:
                logger.warning("Gemini disabled due to dependency conflict")
            return ""

        # ---------- IMAGES ----------
        if ext in [".png", ".jpg", ".jpeg"]:
            if extract_text_with_gemini:
                return await extract_text_with_gemini(file_path, is_pdf=False)
            else:
                logger.warning("Gemini disabled due to dependency conflict")
            return ""

        logger.warning("Unsupported file type: %s", ext)
        return ""

    except Exception as e:
        logger.exception("Text extraction failed: %s", e)
        return ""



try:
    from app.libs.gemini_service import extract_text_with_gemini
except Exception:
    extract_text_with_gemini = None









